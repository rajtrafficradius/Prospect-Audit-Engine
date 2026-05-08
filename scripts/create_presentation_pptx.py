import os
import json
import sys
import subprocess
from datetime import datetime
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION

# Premium executive visual system
BG_COLOR = RGBColor(246, 248, 252)
SURFACE_TINT = RGBColor(250, 252, 255)
TEXT_DARK = RGBColor(13, 27, 66)
TEXT_MID = RGBColor(89, 102, 130)
TEXT_SOFT = RGBColor(142, 151, 171)
ACCENT = RGBColor(30, 94, 255)
ACCENT_2 = RGBColor(80, 135, 255)
ACCENT_3 = RGBColor(201, 220, 255)
ACCENT_DARK = RGBColor(11, 33, 89)
GLASS_FILL = RGBColor(255, 255, 255)
GLASS_BORDER = RGBColor(222, 229, 241)
GRID_COLOR = RGBColor(231, 236, 245)
PANEL_SHADOW = RGBColor(87, 112, 166)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
FONT_FAMILY = "Aptos"
COVER_FONT_FAMILY = "Aptos"
ICON_STROKE = ACCENT

LOGO_LEFT = Cm(25.84)
LOGO_TOP = Cm(0.89)
LOGO_WIDTH = Cm(6.0)
LOGO_HEIGHT = Cm(1.71)
SLIDE_NUM_WIDTH = Cm(2.0)
SLIDE_NUM_HEIGHT = Cm(0.7)
SLIDE_NUM_LEFT = Cm(27.84)
SLIDE_NUM_TOP = Cm(0.22)

# One centered chart area on content slides
CHART_X = Inches(1.0)
CHART_Y = Inches(1.86)
CHART_W = Cm(28.5)
CHART_H = Cm(9.0)

# Bullets below chart
BULLET_X = Inches(0.85)
BULLET_Y = Inches(5.52)
BULLET_W = Cm(30.0)
BULLET_H = Cm(5.1)
CONTENT_BULLET_Y = Inches(5.78)


def _set_bg(slide, prs):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()


def _add_full_slide_background(slide, prs, image_path):
    if not image_path or not os.path.exists(image_path):
        _set_bg(slide, prs)
        return
    try:
        with Image.open(image_path) as img:
            img_w, img_h = img.size
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        scale = max(slide_w / float(img_w), slide_h / float(img_h))
        draw_w = int(round(img_w * scale))
        draw_h = int(round(img_h * scale))
        left = int(round((slide_w - draw_w) / 2))
        top = int(round((slide_h - draw_h) / 2))
        slide.shapes.add_picture(image_path, left, top, width=draw_w, height=draw_h)
    except Exception:
        _set_bg(slide, prs)


def _draw_cover_inspired_background(slide, prs):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(10, 22, 40)
    bg.line.fill.background()

    wash = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.1), Inches(-0.7), Inches(4.2), Inches(4.0))
    wash.fill.solid()
    wash.fill.fore_color.rgb = RGBColor(38, 74, 150)
    wash.fill.transparency = 0.82
    wash.line.fill.background()

    lower = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(4.9), Inches(4.5), Inches(3.7))
    lower.fill.solid()
    lower.fill.fore_color.rgb = RGBColor(18, 46, 104)
    lower.fill.transparency = 0.86
    lower.line.fill.background()

    for left, top, w, h, color, transparency in [
        (Inches(8.7), Inches(4.42), Inches(0.34), Inches(0.95), RGBColor(76, 124, 227), 0.76),
        (Inches(9.18), Inches(4.1), Inches(0.34), Inches(1.28), RGBColor(95, 143, 244), 0.74),
        (Inches(9.66), Inches(3.72), Inches(0.34), Inches(1.66), RGBColor(126, 168, 255), 0.72),
        (Inches(10.14), Inches(3.28), Inches(0.34), Inches(2.1), RGBColor(155, 191, 255), 0.7),
        (Inches(10.62), Inches(2.76), Inches(0.34), Inches(2.62), RGBColor(190, 216, 255), 0.68),
    ]:
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.fill.transparency = transparency
        bar.line.fill.background()

    points = [
        (Inches(1.15), Inches(6.52)),
        (Inches(2.58), Inches(6.08)),
        (Inches(4.08), Inches(5.88)),
        (Inches(5.32), Inches(5.26)),
        (Inches(6.92), Inches(4.85)),
        (Inches(8.3), Inches(4.04)),
        (Inches(9.64), Inches(2.96)),
        (Inches(11.16), Inches(1.76)),
    ]
    for i in range(len(points) - 1):
        line = slide.shapes.add_connector(1, points[i][0], points[i][1], points[i + 1][0], points[i + 1][1])
        line.line.color.rgb = RGBColor(112, 163, 255)
        line.line.width = Pt(1.25)
        line.line.transparency = 0.28
    for x, y in points[1::2]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.08), y - Inches(0.08), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(134, 180, 255)
        dot.fill.transparency = 0.08
        dot.line.fill.background()


def _add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if p.runs:
        run = p.runs[0]
        run.font.name = FONT_FAMILY
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    p.line_spacing = 1.05
    return box


def _add_clean_panel(slide, left, top, width, height, radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    panel = slide.shapes.add_shape(radius_shape, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = GLASS_BORDER
    panel.line.width = Pt(0.9)
    return panel


def _clean_list(items):
    out = []
    for x in items or []:
        t = str(x or "").strip()
        if t:
            out.append(t)
    return out


def _clean_optional_text(value):
    text = "" if value is None else str(value).strip()
    if text.lower() in {"none", "null", "nan"}:
        return ""
    return text


def _short_chart_label(text, max_words=3):
    t = " ".join(str(text or "").split())
    if not t:
        return "Metric"
    words = t.split()
    if len(words) <= max_words:
        return t
    return " ".join(words[:max_words])


def _score_from_text(value_text):
    text = str(value_text or "").strip()
    if not text:
        return 0

    matches = re.findall(r"[-+]?\d[\d,]*\.?\d*", text.replace("$", ""))
    if matches:
        num = float(matches[0].replace(",", ""))
        suffix_match = re.search(re.escape(matches[0]) + r"\s*([KkMm])", text)
        if suffix_match:
            suffix = suffix_match.group(1).lower()
            if suffix == "k":
                num *= 1000
            elif suffix == "m":
                num *= 1000000
        return num

    lowered = text.lower()
    keyword_scores = [
        (["transformative", "leader", "dominant"], 92),
        (["high demand", "high", "strong", "boost", "available"], 78),
        (["current", "moderate", "potential", "capture", "drive"], 62),
        (["needed", "need", "develop", "create", "implement", "enhance", "focus", "critical"], 48),
        (["warning", "inadequate", "limited", "compromised"], 34),
        (["fail", "missing", "absent", "lacking", "weak", "underdeveloped", "none"], 14),
    ]
    for keys, score in keyword_scores:
        if any(k in lowered for k in keys):
            return score
    return 50


def _parse_visual_points(visual_data, max_points=6):
    points = []
    for idx, item in enumerate(_clean_list(visual_data)):
        if ":" in item:
            label, raw_value = item.split(":", 1)
        else:
            label, raw_value = f"Metric {idx + 1}", item
        label = _short_chart_label(label.strip(), max_words=3)
        raw_value = raw_value.strip()
        value = _score_from_text(raw_value)
        points.append({"label": label, "value": value, "raw": raw_value, "has_numeric": bool(re.search(r"\d", raw_value))})
        if len(points) >= max_points:
            break
    return points


def _normalized_points(points, target_max=100):
    vals = [float(p["value"]) for p in points if p.get("value") is not None]
    if not vals:
        return points
    current_max = max(vals) or 1
    if current_max <= target_max:
        return points
    ratio = float(target_max) / current_max
    normalized = []
    for p in points:
        clone = dict(p)
        clone["value"] = round(float(clone["value"]) * ratio, 1)
        normalized.append(clone)
    return normalized


def _prepare_chart_points(points, mode):
    if not points:
        return points

    mode = (mode or "bar").lower()
    if mode in {"bar", "line", "funnel"}:
        numeric_points = [p for p in points if p.get("has_numeric")]
        if len(numeric_points) >= 3:
            return numeric_points[:6]
    return points


def _format_chart_value(value):
    value = float(value or 0)
    if value >= 1000000:
        return f"{value/1000000:.1f}M"
    if value >= 1000:
        return f"{value/1000:.0f}K"
    if abs(value - round(value)) < 0.01:
        return str(int(round(value)))
    return f"{value:.1f}"


def _compact_right_bar_text(value, max_words=2, max_chars=16):
    text = " ".join(str(value or "").split()).strip()
    if not text:
        return ""
    if re.fullmatch(r"[\d,.]+%?", text):
        return text
    words = text.split()
    compact = " ".join(words[:max_words])
    if len(compact) > max_chars:
        compact = compact[:max_chars].rsplit(" ", 1)[0].rstrip(" ,;:-")
    return compact or text[:max_chars].rstrip(" ,;:-")


def _add_icon_badge(slide, left, top, kind, size_cm=0.8):
    badge_size = Cm(size_cm)
    shadow = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Cm(0.02), top + Cm(0.05), badge_size, badge_size)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(188, 199, 220)
    shadow.fill.transparency = 0.86
    shadow.line.fill.background()

    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, badge_size, badge_size)
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(239, 244, 253)
    badge.line.color.rgb = GLASS_BORDER
    badge.line.width = Pt(0.8)

    icon_size = max(0.34, size_cm - 0.28)
    offset = max(0, (size_cm - icon_size) / 2.0)
    _draw_icon(slide, kind, left + Cm(offset), top + Cm(offset), size_cm=icon_size)


def _add_line_summary_banner(slide, points):
    return


def _add_bar_side_labels(slide, points):
    base_left = CHART_X + Cm(0.45)
    base_top = CHART_Y + Cm(0.62)
    row_h = Cm(1.95)
    for i, point in enumerate(points[:4]):
        y = base_top + i * row_h
        _add_icon_badge(slide, base_left, y + Cm(0.22), _icon_kind(point["label"]), size_cm=0.72)
        _add_textbox(slide, base_left + Cm(0.86), y + Cm(0.02), Cm(4.9), Cm(0.5), f"{point['label']}:", size=11.8, bold=True, color=TEXT_DARK)
        _add_textbox(slide, base_left + Cm(0.86), y + Cm(0.5), Cm(5.55), Cm(0.82), point.get("raw") or _format_chart_value(point["value"]), size=9.8, color=TEXT_MID)


def _add_radar_annotations(slide, points):
    positions = [
        (CHART_X + Cm(9.1), CHART_Y - Cm(0.05)),
        (CHART_X + Cm(16.7), CHART_Y + Cm(2.2)),
        (CHART_X + Cm(14.9), CHART_Y + Cm(6.9)),
        (CHART_X + Cm(3.0), CHART_Y + Cm(6.9)),
        (CHART_X + Cm(1.05), CHART_Y + Cm(2.2)),
    ]
    for i, point in enumerate(points[:5]):
        x, y = positions[i]
        _add_icon_badge(slide, x, y, _icon_kind(point["label"]), size_cm=0.62)
        _add_textbox(slide, x + Cm(0.75), y - Cm(0.01), Cm(5.0), Cm(0.42), point["label"], size=10.6, bold=True, color=TEXT_DARK)
        _add_textbox(slide, x + Cm(0.75), y + Cm(0.36), Cm(5.0), Cm(0.36), _format_chart_value(point["value"]), size=10.8, bold=True, color=ACCENT)

    summary_x = CHART_X + Cm(21.2)
    summary_y = CHART_Y + Cm(0.95)
    summary_w = Cm(7.0)
    summary_h = Cm(5.0)
    _add_soft_panel(slide, summary_x, summary_y, summary_w, summary_h)
    _add_textbox(slide, summary_x + Cm(0.45), summary_y + Cm(0.35), Cm(6.0), Cm(0.4), "PERFORMANCE SUMMARY", size=11, bold=True, color=TEXT_DARK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, summary_x + Cm(0.45), summary_y + Cm(0.92), Cm(5.8), Cm(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = GRID_COLOR; line.line.fill.background()
    for idx, point in enumerate(points[:5]):
        row_y = summary_y + Cm(1.28 + idx * 0.67)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, summary_x + Cm(0.45), row_y + Cm(0.05), Cm(0.18), Cm(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT; dot.line.fill.background()
        _add_textbox(slide, summary_x + Cm(0.9), row_y - Cm(0.01), Cm(4.2), Cm(0.28), point["label"], size=10.2, color=TEXT_DARK)
        _add_textbox(slide, summary_x + Cm(5.15), row_y - Cm(0.01), Cm(1.0), Cm(0.28), _format_chart_value(point["value"]), size=10.5, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)


def _add_donut_detail_rows(slide, points, colors):
    legend_left = CHART_X + Cm(0.6)
    legend_top = CHART_Y + Cm(0.82)
    row_h = Cm(1.62)
    for i, point in enumerate(points[:4]):
        y = legend_top + i * row_h
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, legend_left, y + Cm(0.28), Cm(0.22), Cm(0.22))
        dot.fill.solid(); dot.fill.fore_color.rgb = colors[i % len(colors)]; dot.line.fill.background()
        _add_icon_badge(slide, legend_left + Cm(0.42), y + Cm(0.04), _icon_kind(point["label"]), size_cm=0.64)
        _add_textbox(slide, legend_left + Cm(1.2), y + Cm(0.02), Cm(7.1), Cm(0.45), f"{point['label']}: {point.get('raw') or _format_chart_value(point['value'])}", size=11, color=TEXT_DARK)
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_left + Cm(0.42), y + Cm(1.04), Cm(7.25), Cm(0.02))
        sep.fill.solid(); sep.fill.fore_color.rgb = GRID_COLOR; sep.line.fill.background()


def _group_to_max_five_verbatim(items):
    vals = _clean_list(items)
    if len(vals) <= 5:
        return vals
    groups = 5
    chunk = (len(vals) + groups - 1) // groups
    merged = []
    for i in range(0, len(vals), chunk):
        merged.append(" | ".join(vals[i:i + chunk]))
    return merged[:5]


def _to_two_lines_max_verbatim(text, target=74):
    t = " ".join(str(text or "").split())
    if len(t) <= target:
        return t
    split_points = [" | ", "; ", ", ", " - ", " "]
    pivot = target
    best = -1
    best_sep = ""
    for sep in split_points:
        idx = t.rfind(sep, 0, pivot + 1)
        if idx > best:
            best = idx
            best_sep = sep
    if best > 10:
        left = t[:best].rstrip()
        right = t[best + len(best_sep):].lstrip()
        return f"{left}\n{right}"
    return t


def _compact_kpi_text(text, target=74):
    t = " ".join(str(text or "").split())
    if not t:
        return t
    detail_target = max(int(target), 54)
    if len(t) <= detail_target:
        return t

    search_start = max(12, int(detail_target * 0.6))
    search_end = min(len(t), detail_target + 18)
    for sep in ["; ", " | ", ". ", ", ", " - ", ": "]:
        idx = t.rfind(sep, search_start, search_end)
        if idx > 16:
            return t[:idx].rstrip(" ,;:-|")

    return _truncate_words_safely(t, detail_target + 8)


def _sentence_case_fragment(text):
    clean = " ".join(str(text or "").split()).strip(" .;:-|")
    if not clean:
        return ""
    return clean[0].upper() + clean[1:]


def _truncate_words_safely(text, target):
    clean = " ".join(str(text or "").split()).strip()
    if not clean or len(clean) <= target:
        return clean
    cut = clean[:target + 1]
    if " " in cut:
        clean = cut.rsplit(" ", 1)[0]
    else:
        clean = cut[:target]
    return clean.rstrip(" ,;:-")


def _split_long_insight(source):
    clean = " ".join(str(source or "").split()).strip(" .")
    if not clean:
        return []

    lower = clean.lower()
    markers = [", while ", ", with ", ", and ", " while ", " with ", " and ", " but ", " because ", " to "]
    preferred = None
    mid_lo = int(len(clean) * 0.28)
    mid_hi = int(len(clean) * 0.72)
    for marker in markers:
        idx = lower.find(marker)
        if idx != -1 and mid_lo <= idx <= mid_hi:
            preferred = (idx, marker)
            break

    if preferred:
        idx, marker = preferred
        left = clean[:idx].rstrip(" ,;:-")
        right = clean[idx + len(marker):].lstrip(" ,;:-")
        if len(left.split()) >= 4 and len(right.split()) >= 4:
            return [_sentence_case_fragment(left), _sentence_case_fragment(right)]

    words = clean.split()
    if len(words) >= 14:
        split_at = max(6, min(len(words) - 6, len(words) // 2))
        return [
            _sentence_case_fragment(" ".join(words[:split_at])),
            _sentence_case_fragment(" ".join(words[split_at:])),
        ]
    return [_sentence_case_fragment(clean)]


def _contextualize_insight(fragment, context_title="", context_subtitle="", target=44):
    base = _sentence_case_fragment(fragment)
    if not base:
        return ""
    cue = " ".join(str(v or "").strip() for v in [context_title, context_subtitle] if str(v or "").strip())
    if cue and cue.lower() not in base.lower() and len(base.split()) < 6:
        candidate = _truncate_words_safely(f"{cue}: {base.lower()}", target)
        if len(candidate.split()) >= 5:
            return candidate
    return _truncate_words_safely(base, target)


def _kpi_card_points(text, max_points=2, target=44, context_title="", context_subtitle=""):
    source = " ".join(str(text or "").split())
    if not source:
        return []

    max_points = min(max_points, 2)
    detail_target = max(int(target), 54)
    normalized = (
        source.replace(" | ", ". ")
        .replace("; ", ". ")
        .replace(" - ", ". ")
        .replace(" • ", ". ")
    )
    raw_parts = [p.strip(" .") for p in re.split(r"\.\s+|:\s+(?=[A-Z])", normalized) if p.strip(" .")]

    parts = []
    seen = set()
    for part in raw_parts:
        for candidate in _split_long_insight(part):
            compact = _truncate_words_safely(_compact_kpi_text(candidate, target=detail_target + 6), detail_target)
            compact = _to_two_lines_max_verbatim(compact, target=detail_target)
            key = compact.lower()
            if compact and key not in seen:
                parts.append(compact)
                seen.add(key)
            if len(parts) >= max_points:
                break
        if len(parts) >= max_points:
            break

    if not parts:
        parts = [_contextualize_insight(source, context_title, context_subtitle, target=detail_target)]

    if len(parts) == 1:
        expanded = _split_long_insight(source)
        if len(expanded) >= 2:
            candidate = _truncate_words_safely(expanded[1], detail_target)
        else:
            candidate = _contextualize_insight(expanded[0] if expanded else source, context_title, context_subtitle, target=detail_target)
        candidate = _to_two_lines_max_verbatim(candidate, target=detail_target)
        if candidate and candidate.lower() != parts[0].lower():
            parts.append(candidate)

    cleaned = []
    for part in parts[:max_points]:
        compact = _to_two_lines_max_verbatim(_truncate_words_safely(part, detail_target), target=detail_target)
        if compact and len(compact.split()) >= 3:
            cleaned.append(compact)

    return cleaned[:max_points]


def _is_roadmap_slide(slide_data):
    title = str(slide_data.get("title", "")).strip().lower()
    return title == "implementation roadmap & revenue map"


def _strip_phase_prefix(text):
    return re.sub(r"^\s*phase\s*\d+\s*:\s*", "", str(text or ""), flags=re.I).strip()


def _roadmap_short_phrase(text, max_words=3):
    clean = _strip_phase_prefix(text)
    words = clean.split()
    if len(words) <= max_words:
        return clean
    return " ".join(words[:max_words])


def _roadmap_stage_groups(slide_data):
    labels = [_strip_phase_prefix(v) for v in _clean_list(slide_data.get("visual_data", []))]
    details = [_strip_phase_prefix(v) for v in _clean_list(slide_data.get("bullets", []))]

    if not labels:
        labels = details[:]
    if not labels:
        labels = [
            "Technical Foundation",
            "Content Optimization",
            "Authority Enhancement",
            "Growth Expansion",
        ]
    if not details:
        details = labels[:]

    groups = []
    total = max(len(labels), 1)
    for idx in range(4):
        start = int(round(idx * total / 4.0))
        end = int(round((idx + 1) * total / 4.0))
        group_labels = labels[start:end] or [labels[min(idx, len(labels) - 1)]]

        d_total = max(len(details), 1)
        d_start = int(round(idx * d_total / 4.0))
        d_end = int(round((idx + 1) * d_total / 4.0))
        group_details = details[d_start:d_end] or [details[min(idx, len(details) - 1)]]

        descriptor = " + ".join(_roadmap_short_phrase(v, max_words=2) for v in group_labels[:2])
        card_title = _roadmap_short_phrase(group_labels[0], max_words=3)
        detail_source = ". ".join(_clean_list(group_details[:2])) or " ".join(group_details)
        card_note = _compact_kpi_text(detail_source, target=118)
        icon_key = " ".join(group_labels)
        groups.append(
            {
                "phase": f"Phase {idx + 1}",
                "number": f"{idx + 1:02d}",
                "descriptor": descriptor,
                "card_title": card_title,
                "card_note": card_note.rstrip("."),
                "icon": _icon_kind(icon_key),
            }
        )
    return groups


def _draw_roadmap_slide(slide, slide_data):
    _add_title_block(slide, _clean_optional_text(slide_data.get("title", "")), _clean_optional_text(slide_data.get("subtitle", "")))

    panel_left = CHART_X - Cm(0.15)
    panel_top = CHART_Y + Cm(0.05)
    panel_w = CHART_W + Cm(0.3)
    panel_h = Cm(6.45)
    _add_clean_panel(slide, panel_left, panel_top, panel_w, panel_h)

    line_left = panel_left + Cm(1.6)
    line_top = panel_top + Cm(1.75)
    line_w = panel_w - Cm(2.7)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, line_top, line_w, Cm(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(203, 213, 225)
    line.line.fill.background()

    groups = _roadmap_stage_groups(slide_data)
    node_base_left = panel_left + Cm(4.25)
    node_gap = Cm(6.35)
    node_y = line_top - Cm(0.28)
    connector_y = panel_top + Cm(2.18)
    icon_y = panel_top + Cm(2.78)
    label_y = panel_top + Cm(4.0)

    for idx, group in enumerate(groups):
        center_x = node_base_left + idx * node_gap
        inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - Cm(0.4), node_y + Cm(0.04), Cm(0.8), Cm(0.8))
        inner.fill.solid()
        inner.fill.fore_color.rgb = ACCENT
        inner.line.fill.background()
        _add_textbox(slide, center_x - Cm(0.5), node_y + Cm(0.13), Cm(1.0), Cm(0.42), str(idx + 1), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Cm(0.012), connector_y, Cm(0.024), Cm(0.8))
        connector.fill.solid()
        connector.fill.fore_color.rgb = ACCENT_3
        connector.line.fill.background()

        _draw_icon(slide, group["icon"], center_x - Cm(0.28), icon_y, size_cm=0.56)
        _add_textbox(slide, center_x - Cm(2.05), label_y, Cm(4.1), Cm(0.58), group["phase"], size=13.5, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _add_textbox(slide, center_x - Cm(2.45), label_y + Cm(0.66), Cm(4.9), Cm(0.72), group["descriptor"], size=10.5, color=TEXT_MID, align=PP_ALIGN.CENTER)

    card_top = BULLET_Y - Cm(0.12)
    card_h = Cm(3.55)
    card_gap = Cm(0.32)
    card_w = (BULLET_W - Cm(0.18) - (card_gap * 3)) / 4
    for idx, group in enumerate(groups):
        x = BULLET_X + Cm(0.09) + idx * (card_w + card_gap)
        _add_layered_insight_card(slide, x, card_top, card_w, card_h)
        _add_icon_badge(slide, x + Cm(0.5), card_top + Cm(0.48), group["icon"], size_cm=0.82)
        _add_textbox(slide, x + Cm(1.56), card_top + Cm(0.5), card_w - Cm(2.02), Cm(0.5), group["card_title"], size=13.6, bold=True, color=ACCENT_DARK)
        roadmap_points = _kpi_card_points(
            group["card_note"],
            max_points=1,
            target=42,
            context_title=group["card_title"],
            context_subtitle=group["phase"],
        )
        _add_card_bullets(
            slide,
            x + Cm(0.54),
            card_top + Cm(1.18),
            card_w - Cm(0.88),
            card_h - Cm(1.46),
            roadmap_points,
            size=13.6,
            color=TEXT_DARK,
        )


def _stable_hash_text(*parts):
    joined = "||".join(" ".join(str(p or "").split()).lower() for p in parts)
    return sum(ord(ch) for ch in joined)


def _icon_kind(text):
    t = str(text or "").lower()
    if "xml" in t or "sitemap" in t or "schema" in t or "structured" in t:
        return "schema"
    if "meta" in t or "content" in t or "faq" in t or "answer" in t:
        return "content"
    if "authority" in t or "entity" in t or "trust" in t:
        return "authority"
    if "technical" in t or "seo" in t or "health" in t:
        return "technical"
    if "index" in t or "crawl" in t:
        return "indexability"
    if "engagement" in t or "description" in t or "message" in t:
        return "engagement"
    if "optimization" in t or "title" in t:
        return "optimization"
    if "performance" in t:
        return "performance"
    if "growth" in t:
        return "growth"
    if "revenue" in t or "coin" in t or "money" in t:
        return "revenue"
    if "timeline" in t or "phase" in t or "milestone" in t:
        return "timeline"
    if "strategy" in t:
        return "strategy"
    if "traffic" in t or "reach" in t or "visibility" in t:
        return "traffic"
    if "kpi" in t or "value" in t or "score" in t:
        return "kpi"
    return "generic"


def _draw_icon(slide, kind, left, top, size_cm=0.55):
    size = Cm(size_cm)
    line_color = ICON_STROKE
    line_width = Pt(1.1)
    if kind == "schema":
        n1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Cm(0.18), top + Cm(0.02), Cm(0.16), Cm(0.16))
        n2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Cm(0.03), top + Cm(0.36), Cm(0.16), Cm(0.16))
        n3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Cm(0.33), top + Cm(0.36), Cm(0.16), Cm(0.16))
        for n in [n1, n2, n3]:
            n.fill.background(); n.line.color.rgb = line_color; n.line.width = line_width
        l1 = slide.shapes.add_connector(1, left + Cm(0.26), top + Cm(0.18), left + Cm(0.11), top + Cm(0.36))
        l2 = slide.shapes.add_connector(1, left + Cm(0.26), top + Cm(0.18), left + Cm(0.41), top + Cm(0.36))
        for l in [l1, l2]:
            l.line.color.rgb = line_color
            l.line.width = line_width
        return
    if kind == "content":
        doc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Cm(0.08), top + Cm(0.02), Cm(0.38), Cm(0.5))
        doc.fill.background(); doc.line.color.rgb = line_color; doc.line.width = line_width
        for dy in [0.12, 0.22, 0.32]:
            ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Cm(0.16), top + Cm(dy), Cm(0.21), Cm(0.02))
            ln.fill.solid(); ln.fill.fore_color.rgb = line_color; ln.line.fill.background()
        return
    if kind == "authority":
        badge = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left + Cm(0.08), top + Cm(0.02), Cm(0.38), Cm(0.44))
        badge.fill.background(); badge.line.color.rgb = line_color; badge.line.width = line_width
        star = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, left + Cm(0.19), top + Cm(0.12), Cm(0.16), Cm(0.16))
        star.fill.background(); star.line.color.rgb = line_color; star.line.width = Pt(0.9)
        return
    if kind == "technical":
        gear = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left + Cm(0.1), top + Cm(0.08), Cm(0.34), Cm(0.34))
        gear.fill.background(); gear.line.color.rgb = line_color; gear.line.width = line_width
        core = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Cm(0.2), top + Cm(0.18), Cm(0.14), Cm(0.14))
        core.fill.background(); core.line.color.rgb = line_color; core.line.width = Pt(0.9)
        return
    if kind == "indexability":
        t1 = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left + Cm(0.06), top + Cm(0.08), Cm(0.42), Cm(0.28))
        t1.fill.background(); t1.line.color.rgb = line_color; t1.line.width = line_width
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Cm(0.32), top + Cm(0.15), Cm(0.06), Cm(0.06))
        dot.fill.background(); dot.line.color.rgb = line_color; dot.line.width = Pt(0.9)
        return
    if kind == "engagement":
        bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Cm(0.06), top + Cm(0.08), Cm(0.4), Cm(0.28))
        bubble.fill.background(); bubble.line.color.rgb = line_color; bubble.line.width = line_width
        for dx in [0.14, 0.24, 0.34]:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Cm(dx), top + Cm(0.18), Cm(0.04), Cm(0.04))
            dot.fill.solid(); dot.fill.fore_color.rgb = line_color; dot.line.fill.background()
        return
    if kind == "optimization":
        tag = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left + Cm(0.06), top + Cm(0.06), Cm(0.42), Cm(0.3))
        tag.fill.background(); tag.line.color.rgb = line_color; tag.line.width = line_width
        return
    if kind == "performance":
        for i, h in enumerate([0.18, 0.3, 0.42]):
            b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Cm(i * 0.18), top + Cm(0.55 - h), Cm(0.12), Cm(h))
            b.fill.solid(); b.fill.fore_color.rgb = line_color; b.line.fill.background()
        return
    if kind == "growth":
        a = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top + Cm(0.18), Cm(0.48), Cm(0.22))
        a.fill.background(); a.line.color.rgb = line_color; a.line.width = line_width
        return
    if kind == "revenue":
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        c.fill.background(); c.line.color.rgb = line_color; c.line.width = line_width
        _add_textbox(slide, left + Cm(0.16), top + Cm(0.1), Cm(0.22), Cm(0.22), "$", size=10, bold=True, color=line_color, align=PP_ALIGN.CENTER)
        return
    if kind == "timeline":
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Cm(0.05), top + Cm(0.05), Cm(0.45), Cm(0.45))
        c.fill.background(); c.line.color.rgb = line_color; c.line.width = line_width
        h = slide.shapes.add_connector(1, left + Cm(0.27), top + Cm(0.1), left + Cm(0.27), top + Cm(0.26))
        h.line.color.rgb = line_color
        h.line.width = line_width
        return
    if kind == "strategy":
        t = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, size, size)
        t.fill.background(); t.line.color.rgb = line_color; t.line.width = line_width
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Cm(0.2), top + Cm(0.2), Cm(0.15), Cm(0.15))
        d.fill.background(); d.line.color.rgb = line_color; d.line.width = Pt(0.9)
        return
    if kind == "traffic":
        for dx, dy in [(0.05, 0.25), (0.25, 0.05), (0.45, 0.25)]:
            n = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Cm(dx), top + Cm(dy), Cm(0.1), Cm(0.1))
            n.fill.background(); n.line.color.rgb = line_color; n.line.width = Pt(0.9)
        return
    if kind == "kpi":
        g = slide.shapes.add_shape(MSO_SHAPE.PIE, left, top, size, size)
        g.fill.background()
        g.line.color.rgb = line_color
        g.line.width = line_width
        return
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    d.fill.background(); d.line.color.rgb = line_color; d.line.width = line_width


def _apply_glass(shape, transparency=0.18):
    shape.fill.solid()
    shape.fill.fore_color.rgb = GLASS_FILL
    shape.fill.transparency = transparency
    shape.line.color.rgb = GLASS_BORDER
    shape.line.width = Pt(0.8)


def _add_soft_panel(slide, left, top, width, height, radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shadow = slide.shapes.add_shape(radius_shape, left + Cm(0.06), top + Cm(0.12), width, height)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(190, 203, 226)
    shadow.fill.transparency = 0.9
    shadow.line.fill.background()

    panel = slide.shapes.add_shape(radius_shape, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = GLASS_BORDER
    panel.line.width = Pt(0.9)
    return panel


def _add_layered_insight_card(slide, left, top, width, height):
    back = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Cm(0.08), top - Cm(0.08), width, height)
    back.fill.solid()
    back.fill.fore_color.rgb = RGBColor(220, 229, 244)
    back.fill.transparency = 0.42
    back.line.color.rgb = RGBColor(204, 216, 236)
    back.line.width = Pt(0.55)

    front = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    front.fill.solid()
    front.fill.fore_color.rgb = WHITE
    front.line.color.rgb = GLASS_BORDER
    front.line.width = Pt(0.85)
    return front


def _iter_emphasis_segments(text):
    value = str(text or "")
    if not value:
        return []

    pattern = re.compile(
        r"(\b(?:SEO|AEO|GEO|CRO|Schema|Content|Authority|Traffic|Revenue|Visibility|Conversions?|Featured Snippets?|PAA|Market Value|Organic Traffic|Domain strength)\b|\b[A-Z][A-Za-z]+(?:\s+[A-Z][A-Za-z]+){1,2}\b)"
    )
    segments = []
    cursor = 0
    for match in pattern.finditer(value):
        if match.start() > cursor:
            segments.append((value[cursor:match.start()], False))
        segments.append((match.group(0), True))
        cursor = match.end()
    if cursor < len(value):
        segments.append((value[cursor:], False))
    return segments or [(value, False)]


def _add_card_bullets(slide, left, top, width, height, items, size=11.2, color=TEXT_DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    bullet_items = [str(i or "").strip() for i in items if str(i or "").strip()]
    for idx, item in enumerate(bullet_items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(5)
        p.line_spacing = 1.16

        bullet_run = p.add_run()
        bullet_run.text = "•  "
        bullet_run.font.name = FONT_FAMILY
        bullet_run.font.size = Pt(size)
        bullet_run.font.bold = False
        bullet_run.font.underline = False
        bullet_run.font.color.rgb = ACCENT_DARK

        for seg_text, seg_bold in _iter_emphasis_segments(item):
            run = p.add_run()
            run.text = seg_text
            run.font.name = FONT_FAMILY
            run.font.size = Pt(size)
            run.font.bold = seg_bold
            run.font.underline = False
            run.font.color.rgb = color
    return box


def _style_axis(axis, grid=False):
    try:
        axis.format.line.color.rgb = GLASS_BORDER
        axis.format.line.width = Pt(0.8)
    except Exception:
        pass
    try:
        axis.tick_labels.font.name = FONT_FAMILY
        axis.tick_labels.font.size = Pt(10.5)
        axis.tick_labels.font.color.rgb = TEXT_MID
    except Exception:
        pass
    try:
        axis.has_major_gridlines = bool(grid)
        if grid:
            axis.major_gridlines.format.line.color.rgb = GRID_COLOR
            axis.major_gridlines.format.line.width = Pt(0.8)
    except Exception:
        pass


def _style_chart_common(chart, legend_position=None):
    chart.has_title = False
    try:
        chart.chart_area.format.fill.background()
        chart.chart_area.format.line.fill.background()
    except Exception:
        pass
    try:
        chart.plot_area.format.fill.background()
        chart.plot_area.format.line.fill.background()
    except Exception:
        pass
    try:
        if chart.has_legend:
            chart.legend.position = legend_position or XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.name = FONT_FAMILY
            chart.legend.font.size = Pt(10)
            chart.legend.font.color.rgb = TEXT_MID
    except Exception:
        pass

    try:
        _style_axis(chart.category_axis, grid=False)
    except Exception:
        pass
    try:
        _style_axis(chart.value_axis, grid=True)
        chart.value_axis.minimum_scale = 0
    except Exception:
        pass


def _add_chart_panel(slide):
    panel_left = CHART_X - Cm(0.35)
    panel_top = CHART_Y - Cm(0.28)
    panel_w = CHART_W + Cm(0.7)
    panel_h = CHART_H + Cm(0.6)
    _add_clean_panel(slide, panel_left, panel_top, panel_w, panel_h)


def _add_slide2_kpi_panel(slide):
    panel_left = CHART_X - Cm(0.18)
    panel_top = CHART_Y + Cm(0.08)
    panel_w = CHART_W + Cm(0.36)
    panel_h = Cm(5.85)
    _add_clean_panel(slide, panel_left, panel_top, panel_w, panel_h)
    trim = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, panel_left + Cm(0.42), panel_top + Cm(0.24), panel_w - Cm(0.84), Cm(0.07))
    trim.fill.solid()
    trim.fill.fore_color.rgb = ACCENT_3
    trim.fill.transparency = 0.25
    trim.line.fill.background()


def _slide2_kpi_panel_bottom():
    return CHART_Y + Cm(0.08) + Cm(5.85)


def _add_slide2_metric_cards(slide, metrics):
    semrush_unavailable = bool(metrics.get("semrush_unavailable"))
    def _coerce_number(raw):
        try:
            if raw in (None, ""):
                return None
            return float(raw)
        except Exception:
            return None

    labels = [
        (
            "Organic Traffic",
            metrics.get("organic_traffic"),
            "SEMrush unavailable for this run" if semrush_unavailable else "Current organic baseline",
        ),
        (
            "Authority Score",
            metrics.get("authority_score"),
            "SEMrush unavailable for this run" if semrush_unavailable else "Domain strength indicator",
        ),
        (
            "Market Value",
            metrics.get("market_value"),
            "SEMrush unavailable for this run" if semrush_unavailable else "Estimated traffic value",
        ),
    ]

    card_top = CHART_Y + Cm(0.56)
    card_h = Cm(4.86)
    card_w = Cm(8.76)
    gap = Cm(0.54)
    start_left = CHART_X + Cm(0.47)
    accents = [ACCENT, ACCENT_DARK, ACCENT_2]

    numeric_values = [val for _, v, _ in labels if (val := _coerce_number(v)) is not None]
    max_value = max(numeric_values) if numeric_values else 1
    for idx, (label, value, note) in enumerate(labels):
        left = start_left + idx * (card_w + gap)
        _add_soft_panel(slide, left, card_top, card_w, card_h)

        strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Cm(0.42), card_top + Cm(0.26), card_w - Cm(0.84), Cm(0.11))
        strip.fill.solid()
        strip.fill.fore_color.rgb = accents[idx % len(accents)]
        strip.line.fill.background()

        _add_icon_badge(slide, left + Cm(0.56), card_top + Cm(0.76), _icon_kind(label), size_cm=1.08)
        _add_textbox(slide, left + Cm(1.96), card_top + Cm(0.83), card_w - Cm(3.02), Cm(0.38), label, size=11.3, bold=True, color=TEXT_DARK)
        numeric_value = _coerce_number(value)
        display_value = f"{numeric_value:,.0f}" if numeric_value is not None else "N/A"
        value_color = ACCENT if display_value != "N/A" else TEXT_MID
        _add_textbox(slide, left + Cm(1.96), card_top + Cm(1.43), card_w - Cm(3.02), Cm(0.98), display_value, size=30, bold=True, color=value_color)
        _add_textbox(slide, left + Cm(1.96), card_top + Cm(2.5), card_w - Cm(3.02), Cm(0.46), note, size=10, color=TEXT_MID)

        spark_base = left + card_w - Cm(1.68)
        spark_y = card_top + Cm(1.14)
        spark_heights = [0.1, 0.2, 0.32, 0.46, 0.62, 0.8] if display_value != "N/A" else [0.18, 0.18, 0.18, 0.18, 0.18, 0.18]
        spark_color = ACCENT_3 if display_value != "N/A" else GRID_COLOR
        for s_idx, h in enumerate(spark_heights):
            bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, spark_base + Cm(0.17 * s_idx), spark_y + Cm(0.82 - h), Cm(0.1), Cm(h))
            bar.fill.solid()
            bar.fill.fore_color.rgb = spark_color
            bar.line.fill.background()

        track_left = left + Cm(0.52)
        track_top = card_top + card_h - Cm(0.6)
        track_w = card_w - Cm(1.04)
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, track_left, track_top, track_w, Cm(0.16))
        track.fill.solid()
        track.fill.fore_color.rgb = GRID_COLOR
        track.line.fill.background()

        if numeric_value is not None:
            fill_w = int(max(Cm(0.72), track_w * (numeric_value / max_value)))
            fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, track_left, track_top, fill_w, Cm(0.16))
            fill.fill.solid()
            fill.fill.fore_color.rgb = accents[idx % len(accents)]
            fill.line.fill.background()


def _collect_existing_slide_points(slide_data):
    points = []
    points.extend(_clean_list(slide_data.get("bullets", [])))
    takeaway = str(slide_data.get("takeaway", "")).strip()
    if takeaway:
        points.append(takeaway)
    points.extend(_clean_list(slide_data.get("visual_data", [])))
    subtitle = _clean_optional_text(slide_data.get("subtitle", ""))
    if subtitle and len(points) < 3:
        points.append(subtitle)
    title = _clean_optional_text(slide_data.get("title", ""))
    if title and not points:
        points.append(title)

    seen = set()
    cleaned = []
    for point in points:
        normalized = " ".join(str(point or "").split())
        if not normalized:
            continue
        key = normalized.lower()
        if key in seen:
            continue
        seen.add(key)
        cleaned.append(normalized)
    return cleaned


def _resolve_logo_path(project_root):
    candidates = []
    env_logo = os.environ.get("PPT_LOGO_PATH", "").strip()
    if env_logo:
        candidates.append(env_logo)
    candidates.append(os.path.join(project_root, "src", "static", "logo.png"))
    candidates.append(os.path.join(project_root, "templates", "traffikis_logo.png"))
    for c in candidates:
        if c and os.path.exists(c):
            return c
    return ""

def _resolve_cover_background_path(project_root):
    candidates = []
    env_bg = os.environ.get("PPT_COVER_BG_PATH", "").strip()
    if env_bg:
        candidates.append(env_bg)
    candidates.append(os.path.join(project_root, "src", "static", "ppt_assets", "cover_bg.png"))
    candidates.append(os.path.join(project_root, "src", "static", "ppt_assets", "bg_premium.png"))
    candidates.append(os.path.join(project_root, "front_slide.png"))
    candidates.append(os.path.join(project_root, "templates", "front_slide.png"))
    for c in candidates:
        if c and os.path.exists(c):
            return c
    return ""


def _resolve_white_logo_path(project_root):
    candidates = [
        os.path.join(project_root, "src", "static", "logo-white.png"),
        os.path.join(project_root, "src", "static", "logo_white.png"),
    ]
    for c in candidates:
        if c and os.path.exists(c):
            return c
    return _resolve_logo_path(project_root)


def _resolve_content_background_path(project_root):
    candidates = []
    env_bg = os.environ.get("PPT_CONTENT_BG_PATH", "").strip()
    if env_bg:
        candidates.append(env_bg)
    candidates.append(os.path.join(project_root, "src", "static", "ppt_assets", "content_bg.png"))
    candidates.append(os.path.join(project_root, "src", "static", "ppt_assets", "bg_premium.png"))
    candidates.append(os.path.join(project_root, "templates", "content_slide_bg.png"))
    for c in candidates:
        if c and os.path.exists(c):
            return c
    return ""


def _normalize_to_month_year(value):
    text = " ".join(str(value or "").replace(",", " ").split())
    if not text:
        return ""

    for fmt in ("%B %d %Y", "%b %d %Y", "%d %B %Y", "%d %b %Y", "%Y-%m-%d", "%d-%m-%Y", "%m-%d-%Y", "%Y/%m/%d", "%d/%m/%Y", "%m/%d/%Y"):
        try:
            return datetime.strptime(text, fmt).strftime("%B %Y")
        except ValueError:
            pass

    parts = text.split()
    if len(parts) >= 2 and parts[-1].isdigit() and len(parts[-1]) == 4:
        for fmt in ("%B %Y", "%b %Y"):
            try:
                return datetime.strptime(f"{parts[0]} {parts[-1]}", fmt).strftime("%B %Y")
            except ValueError:
                pass

    return text

def _load_cover_context(session_dir):
    context = {}
    for fname in ["business_analysis.json", "strategy_narrative.json", "audit_findings.json"]:
        path = os.path.join(session_dir, fname)
        if os.path.exists(path):
            try:
                with open(path, "r", encoding="utf-8-sig") as f:
                    context[fname] = json.load(f)
            except Exception:
                context[fname] = {}
        else:
            context[fname] = {}
    return context

def _load_market_metrics(session_dir):
    path = os.path.join(session_dir, "market_intelligence.json")
    try:
        with open(path, "r", encoding="utf-8-sig") as f:
            data = json.load(f)
    except Exception:
        data = {}
    prospect = data.get("prospect", {}) or {}
    semrush_unavailable = str(((data.get("availability", {}) or {}).get("semrush") or "")).lower() == "unavailable"
    return {
        "organic_traffic": None if semrush_unavailable else prospect.get("organic_traffic"),
        "authority_score": None if semrush_unavailable else prospect.get("authority_score"),
        "market_value": None if semrush_unavailable else prospect.get("organic_traffic_value"),
        "semrush_unavailable": semrush_unavailable,
        "warning": ((data.get("availability", {}) or {}).get("warning") or ""),
    }

def _build_dynamic_subtitle(context):
    business = context.get("business_analysis.json", {}) or {}
    narrative = context.get("strategy_narrative.json", {}) or {}

    industry = (
        business.get("industry")
        or business.get("company_profile", {}).get("industry")
        or narrative.get("industry_context")
        or "Digital Services"
    )

    challenge = (
        business.get("primary_challenge")
        or narrative.get("executive_summary")
        or narrative.get("company_profile")
        or "improving market visibility and conversion outcomes"
    )
    challenge_text = " ".join(str(challenge).split())
    first_sentence = challenge_text.split(".")[0].strip()
    if not first_sentence:
        first_sentence = "improving market visibility and conversion outcomes"
    first_sentence = first_sentence.rstrip(" ,;:-")
    return f"For {industry}, the primary focus is {first_sentence}."

def _safe_text(value, fallback="N/A"):
    t = " ".join(str(value or "").split())
    return t if t else fallback

def _short_phrase(value, max_words=5):
    t = " ".join(str(value or "").split())
    if not t:
        return "N/A"
    words = t.split()
    return " ".join(words[:max_words])

def _build_cover_profile_items(context):
    business = context.get("business_analysis.json", {}) or {}
    narrative = context.get("strategy_narrative.json", {}) or {}
    findings = context.get("audit_findings.json", {}) or {}

    company_type = (
        business.get("business_type")
        or business.get("company_profile", {}).get("business_type")
        or narrative.get("company_profile")
    )

    primary_focus = (
        business.get("commercial_intent_focus")
        or business.get("description")
        or (business.get("primary_services") or [None])[0]
    )
    primary_focus = _short_phrase(primary_focus, max_words=5)

    strategy_parts = []
    if findings.get("seo_findings"):
        strategy_parts.append("SEO")
    if findings.get("aeo_findings"):
        strategy_parts.append("AEO")
    if findings.get("geo_findings"):
        strategy_parts.append("GEO")
    if not strategy_parts:
        strategy_parts = ["SEO", "AEO", "GEO"]
    digital_strategy = " + ".join(strategy_parts)

    services_count = len(business.get("primary_services") or [])
    money_pages_count = len(business.get("money_pages") or [])
    has_testimonials = bool(business.get("has_testimonials"))
    has_contact_form = bool(business.get("has_contact_form"))
    assets = (
        f"Services: {services_count} | Money Pages: {money_pages_count} | "
        f"Testimonials: {'Yes' if has_testimonials else 'No'} | Contact Form: {'Yes' if has_contact_form else 'No'}"
    )

    return [
        ("Company Type", _safe_text(company_type)),
        ("Primary Focus", _safe_text(primary_focus)),
        ("Digital Strategy", _safe_text(digital_strategy)),
        ("Assets", _safe_text(assets)),
    ]


def _cover_textbox(slide, left, top, width, height, text, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name=COVER_FONT_FAMILY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if p.runs:
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    p.line_spacing = 1.0
    return box

def _add_logo(slide, logo_path):
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, LOGO_LEFT, LOGO_TOP, width=LOGO_WIDTH, height=LOGO_HEIGHT)


def _add_header(slide, index, total, logo_path, include_logo=True):
    _add_textbox(slide, Inches(0.62), Inches(0.17), Inches(3.2), Inches(0.25), "STRATEGIC MILESTONE", size=10, bold=True, color=ACCENT)

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, SLIDE_NUM_LEFT, SLIDE_NUM_TOP, SLIDE_NUM_WIDTH, SLIDE_NUM_HEIGHT)
    pill.fill.solid()
    pill.fill.fore_color.rgb = WHITE
    pill.line.color.rgb = GLASS_BORDER
    pill.line.width = Pt(0.8)
    _add_textbox(slide, SLIDE_NUM_LEFT, SLIDE_NUM_TOP + Cm(0.08), SLIDE_NUM_WIDTH, Cm(0.3), f"{index:02d}/{total:02d}", size=11, bold=False, color=RGBColor(55, 65, 81), align=PP_ALIGN.CENTER)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SLIDE_NUM_LEFT + SLIDE_NUM_WIDTH + Cm(0.38), SLIDE_NUM_TOP + Cm(0.08), Cm(0.02), LOGO_HEIGHT - Cm(0.18))
    sep.fill.solid()
    sep.fill.fore_color.rgb = GRID_COLOR
    sep.line.fill.background()
    if include_logo:
        _add_logo(slide, logo_path)


def _add_title_block(slide, title, subtitle):
    title = _clean_optional_text(title)
    subtitle = _clean_optional_text(subtitle)
    title_left = Cm(1.84)
    title_top = Cm(1.06)
    title_w = Cm(23.6)
    size = 28
    h = Cm(1.48)
    sub_left = title_left
    sub_top = Cm(2.18)
    _add_textbox(slide, title_left, title_top, title_w, h, title, size=size, bold=True, color=TEXT_DARK)
    if subtitle:
        _add_textbox(slide, sub_left, sub_top, Cm(24.8), Cm(0.62), subtitle, size=18, color=TEXT_MID)


def _resize_test_box_5(slide):
    target_w = Cm(26.12)
    target_h = Cm(1.82)
    for shape in slide.shapes:
        if getattr(shape, "name", "").strip().lower() == "textbox 5":
            shape.width = target_w
            shape.height = target_h


def _draw_bar_chart(slide, points):
    rows = points[:4] if points else [{"label": x, "value": v, "raw": str(v)} for x, v in zip(["A", "B", "C", "D"], [85, 72, 63, 54])]
    _draw_horizontal_bar_panel(slide, rows, title_text="PERFORMANCE SUMMARY")


def _draw_pie_chart(slide, points):
    points = points[:4] if points else [{"label": x, "value": v} for x, v in zip(["A", "B", "C", "D"], [30, 24, 20, 16])]
    cats = [p["label"] for p in points]
    vals = [max(1, float(p["value"])) for p in points]
    pie_colors = [
        ACCENT_DARK,
        ACCENT,
        ACCENT_2,
        RGBColor(122, 155, 255),
        RGBColor(188, 205, 236),
    ]
    data = CategoryChartData()
    data.categories = cats
    data.add_series("Distribution", vals)
    # Match reference: pie on right, custom legend block beside it on left.
    chart_left = CHART_X + Cm(11.1)
    chart_top = CHART_Y + Cm(0.35)
    chart_w = Cm(15.3)
    chart_h = Cm(7.5)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, chart_left, chart_top, chart_w, chart_h, data)
    chart = chart_shape.chart
    chart.has_legend = False
    _style_chart_common(chart)
    for i, pt in enumerate(chart.series[0].points):
        pt.format.fill.solid()
        color = pie_colors[i % len(pie_colors)]
        pt.format.fill.fore_color.rgb = color
        pt.format.line.fill.background()

    _add_donut_detail_rows(slide, points, pie_colors)


def _draw_line_chart(slide, points):
    points = points[:6] if points else [{"label": x, "value": v} for x, v in zip(["P1", "P2", "P3", "P4", "P5", "P6"], [42, 56, 51, 68, 64, 79])]
    cats = [p["label"] for p in points]
    vals = [float(p["value"]) for p in points]
    data = CategoryChartData()
    data.categories = cats
    data.add_series("Trend", vals)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, CHART_X + Cm(0.1), CHART_Y + Cm(0.12), CHART_W - Cm(0.2), CHART_H - Cm(1.55), data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    _style_chart_common(chart)
    s = chart.series[0]
    s.format.line.color.rgb = ACCENT
    s.format.line.width = Pt(2.5)
    s.marker.style = 2
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.position = XL_LABEL_POSITION.ABOVE
    plot.data_labels.font.name = FONT_FAMILY
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.color.rgb = TEXT_MID
    _add_line_summary_banner(slide, points)


def _draw_funnel_chart(slide, points):
    points = points[:4] if points else [{"label": x, "value": v} for x, v in zip(["Stage 1", "Stage 2", "Stage 3", "Stage 4"], [100, 72, 51, 34])]
    cats = [p["label"] for p in points]
    vals = [float(p["value"]) for p in points]
    data = CategoryChartData()
    data.categories = cats
    data.add_series("Funnel", vals)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, CHART_X + Cm(0.2), CHART_Y + Cm(0.18), CHART_W - Cm(0.4), CHART_H - Cm(0.36), data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    _style_chart_common(chart)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.name = FONT_FAMILY
    plot.data_labels.font.size = Pt(11)
    plot.data_labels.font.color.rgb = TEXT_DARK
    for i, pt in enumerate(chart.series[0].points):
        pt.format.fill.solid()
        color = [ACCENT_DARK, ACCENT, ACCENT_2, RGBColor(155, 182, 255)][i % 4]
        pt.format.fill.fore_color.rgb = color


def _draw_radar_chart(slide, points):
    points = _normalized_points(points[:5]) if points else [{"label": x, "value": v} for x, v in zip(["A", "B", "C", "D", "E"], [64, 79, 58, 71, 66])]
    cats = [p["label"] for p in points]
    if len(cats) < 5:
        last = cats[-1] if cats else "Metric"
        cats = cats + [last] * (5 - len(cats))
    vals = [float(p["value"]) for p in points]
    if len(vals) < 5:
        vals = vals + [vals[-1] if vals else 50] * (5 - len(vals))
    data = CategoryChartData()
    data.categories = cats
    data.add_series("Coverage", vals)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR_MARKERS, CHART_X + Cm(4.1), CHART_Y + Cm(0.55), Cm(14.4), CHART_H - Cm(1.15), data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    _style_chart_common(chart)
    s = chart.series[0]
    s.format.line.color.rgb = ACCENT
    s.format.line.width = Pt(2.2)
    try:
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = ACCENT_3
        s.format.fill.transparency = 0.58
    except Exception:
        pass
    for i, pt in enumerate(s.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = ACCENT_2
    _add_radar_annotations(slide, points)


def _draw_scatter_chart(slide, points):
    data = XyChartData()
    series = data.add_series("Points")
    points = _normalized_points(points[:6], target_max=10) if points else []
    pts = []
    if points:
        for idx, p in enumerate(points, start=1):
            pts.append((float(idx) * 1.5, max(1.2, float(p["value"]) / 2.0)))
    else:
        pts = [(1.8, 4.5), (3.0, 3.4), (4.4, 3.9), (5.9, 2.9), (7.5, 3.6), (9.7, 2.7)]
    for x, y in pts:
        series.add_data_point(x, y)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, CHART_X + Cm(0.15), CHART_Y + Cm(0.15), CHART_W - Cm(0.3), CHART_H - Cm(0.3), data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    _style_chart_common(chart)
    chart.series[0].format.line.fill.background()
    try:
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = ACCENT
    except Exception:
        pass
    try:
        chart.series[0].marker.style = 2
    except Exception:
        pass


def _draw_progress_metric_slide(slide, points):
    rows = points[:4] if points else [{"label": x, "value": v} for x, v in zip(["SEO", "AEO", "GEO", "Content"], [64, 79, 58, 71])]
    _draw_horizontal_bar_panel(slide, rows, title_text="PERFORMANCE SUMMARY")


def _draw_horizontal_bar_panel(slide, rows, title_text="PERFORMANCE SUMMARY"):
    rows = rows[:4]
    panel_left = CHART_X + Cm(0.55)
    panel_top = CHART_Y + Cm(0.34)
    panel_w = CHART_W - Cm(1.1)
    panel_h = CHART_H - Cm(0.55)
    _add_clean_panel(slide, panel_left, panel_top, panel_w, panel_h)

    _add_textbox(slide, panel_left + Cm(0.62), panel_top + Cm(0.38), Cm(6.0), Cm(0.4), title_text, size=9, bold=True, color=TEXT_MID)

    label_left = panel_left + Cm(0.62)
    label_w = Cm(5.6)
    value_w = Cm(2.25)
    track_left = label_left + label_w
    right_pad = Cm(0.78)
    track_w = panel_w - (track_left - panel_left) - value_w - right_pad - Cm(0.28)
    value_left = panel_left + panel_w - value_w - right_pad

    row_top = panel_top + Cm(1.28)
    row_gap = Cm(1.55)
    bar_h = Cm(0.16)

    palette = [ACCENT, ACCENT_DARK, ACCENT_2, RGBColor(122, 155, 255)]

    for idx, point in enumerate(rows):
        y = row_top + idx * row_gap
        val = max(0.0, min(100.0, float(point.get("value", 0))))
        label = str(point.get("label", f"Metric {idx+1}"))
        raw_value = point.get("raw") or f"{int(round(val))}"
        value_text = _compact_right_bar_text(raw_value)

        _add_textbox(slide, label_left, y - Cm(0.02), label_w - Cm(0.18), Cm(0.34), label, size=11, bold=True, color=TEXT_DARK)
        _add_textbox(slide, value_left, y - Cm(0.03), value_w, Cm(0.38), value_text, size=10.4, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, track_left, y + Cm(0.5), track_w, bar_h)
        track.fill.solid()
        track.fill.fore_color.rgb = GRID_COLOR
        track.line.fill.background()

        fill_w = max(Cm(0.42), track_w * (val / 100.0))
        fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, track_left, y + Cm(0.5), fill_w, bar_h)
        fill.fill.solid()
        fill.fill.fore_color.rgb = palette[idx % len(palette)]
        fill.line.fill.background()


def _draw_center_chart(slide, mode, points):
    m = (mode or "bar").lower()
    if m == "progress":
        _draw_progress_metric_slide(slide, points)
    elif m == "pie":
        _draw_pie_chart(slide, points)
    elif m == "funnel":
        _draw_funnel_chart(slide, points)
    elif m == "line":
        _draw_line_chart(slide, points)
    elif m == "radar":
        _draw_radar_chart(slide, points)
    else:
        _draw_bar_chart(slide, points)


def _draw_slide2_feature_cards(slide, points):
    pts = _group_to_max_five_verbatim(points)
    if len(pts) > 3:
        pts = pts[:3]
    while len(pts) < 3:
        pts.append("")

    card_top = _slide2_kpi_panel_bottom() + Cm(0.42)
    card_h = Cm(4.42)
    gap = Cm(0.42)
    card_w = (BULLET_W - Cm(0.18) - (gap * 2)) / 3
    starts = [BULLET_X + Cm(0.09) + i * (card_w + gap) for i in range(3)]
    titles = ["Current Position", "Strategic Priority", "Growth Outcome"]

    for idx, point in enumerate(pts[:3]):
        x = starts[idx]
        _add_layered_insight_card(slide, x, card_top, card_w, card_h)
        _add_icon_badge(slide, x + Cm(0.52), card_top + Cm(0.44), _icon_kind(point or titles[idx]), size_cm=0.84)
        _add_textbox(slide, x + Cm(1.82), card_top + Cm(0.52), card_w - Cm(2.34), Cm(0.42), titles[idx], size=13.3, bold=True, color=TEXT_DARK)
        card_points = _kpi_card_points(
            point,
            max_points=2,
            target=52,
            context_title=titles[idx],
            context_subtitle="Executive strategy summary",
        )
        _add_card_bullets(
            slide,
            x + Cm(0.88),
            card_top + Cm(1.34),
            card_w - Cm(1.38),
            card_h - Cm(1.72),
            card_points,
            size=14.0,
            color=RGBColor(55, 65, 81),
        )

def _draw_slide2_company_metrics(slide, metrics):
    _add_slide2_metric_cards(slide, metrics)


def _chart_mode_preferences(slide_data, index):
    visual_type = str(slide_data.get("visual_type", "")).strip().lower()
    title = _clean_optional_text(slide_data.get("title", "")).lower()
    subtitle = _clean_optional_text(slide_data.get("subtitle", "")).lower()
    visual_data = _clean_list(slide_data.get("visual_data", []))
    data_text = " ".join(visual_data).lower()
    points = _parse_visual_points(visual_data, max_points=6)
    numeric_points = [p for p in points if p.get("has_numeric")]

    prefs = []

    def add(mode):
        if mode not in prefs:
            prefs.append(mode)

    if "keyword landscape" in subtitle or "economic value" in title:
        add("line")
        add("bar")
    if "geo & entity" in subtitle or "authority" in title:
        add("pie")
        add("radar")
    if "projected impact" in title:
        add("radar")
        add("bar")
    if "conversion engine" in title:
        add("funnel")
        add("bar")
    if "technical evidence" in title and "geo" in subtitle:
        add("pie")
        add("bar")
    if "technical evidence" in title and ("seo" in subtitle or "aeo" in subtitle):
        add("bar")
        add("progress")
    if any(k in title or k in subtitle for k in ["trend", "growth", "timeline", "forecast", "roadmap"]):
        add("line")
        add("bar")
    if any(k in title or k in subtitle for k in ["mix", "distribution", "share", "gap split", "layer distribution"]):
        add("pie")
        add("bar")
    if any(k in title or k in subtitle for k in ["matrix", "opportunity matrix", "cluster", "positioning"]):
        add("bar")
        add("pie")
    if any(k in title or k in subtitle for k in ["readiness", "scorecard", "audit score", "benchmark"]):
        add("radar")
        add("progress")
    if any(k in title or k in subtitle for k in ["path", "conversion", "journey", "stage", "funnel"]):
        add("funnel")
        add("line")

    visual_map = {
        "radar": ["radar", "progress"],
        "comparison": ["bar", "progress"],
        "matrix": ["bar", "pie"],
        "funnel": ["funnel", "bar"],
        "architecture": ["line", "bar"],
        "pyramid": ["line", "bar"],
        "distribution": ["pie", "bar"],
        "trend": ["line", "bar"],
        "scorecard": ["radar", "progress"],
        "scatter": ["bar", "line"],
    }
    if visual_type in visual_map:
        for mode in visual_map[visual_type]:
            add(mode)

    if len(numeric_points) >= 4:
        add("bar")
        add("radar")
    elif len(numeric_points) >= 3:
        add("line")
        add("bar")
    elif len(points) >= 3:
        add("pie")
        add("bar")

    signature = _stable_hash_text(title, subtitle, data_text, visual_type, index)
    fallback_modes = ["bar", "line", "pie", "radar", "funnel", "progress"]
    start = signature % len(fallback_modes)
    rotated = fallback_modes[start:] + fallback_modes[:start]
    for mode in rotated:
        add(mode)
    return prefs


def _select_chart_mode(slide_data, index, recent_modes=None):
    recent_modes = list(recent_modes or [])
    prefs = _chart_mode_preferences(slide_data, index)
    if not prefs:
        return "bar"

    last_mode = recent_modes[-1] if recent_modes else None
    repeated_pair = len(recent_modes) >= 2 and recent_modes[-1] == recent_modes[-2]

    for mode in prefs:
        if last_mode and mode == last_mode and len(prefs) > 1:
            continue
        if repeated_pair and mode == recent_modes[-1] and len(prefs) > 2:
            continue
        return mode

    return prefs[0]


def _add_bottom_bullets(slide, points, slide_index, slide_title="", slide_subtitle=""):
    if slide_index == 2:
        _draw_slide2_feature_cards(slide, points)
        return
    pts = _group_to_max_five_verbatim(points)
    if len(pts) > 3:
        pts = pts[:3]
    if len(pts) < 2:
        return

    display_pts = pts[:3]
    rows = 1
    cols = len(display_pts)
    card_h = Cm(3.8)
    row_gap = Cm(0.0)
    col_gap = Cm(0.56)
    col_w = (BULLET_W - Cm(0.18) - (col_gap * (cols - 1))) / cols
    for i, point in enumerate(display_pts):
        row = i // cols
        col = i % cols
        x = BULLET_X + Cm(0.09) + col * (col_w + col_gap)
        y = CONTENT_BULLET_Y + row * (card_h + row_gap)
        card_w = col_w
        _add_layered_insight_card(slide, x, y, card_w, card_h)
        _add_icon_badge(slide, x + Cm(0.46), y + Cm(0.48), _icon_kind(point), size_cm=0.78)
        card_points = _kpi_card_points(
            point,
            max_points=2,
            target=42 if cols >= 3 else 48,
            context_title=slide_title,
            context_subtitle=slide_subtitle,
        )
        _add_card_bullets(
            slide,
            x + Cm(1.46),
            y + Cm(0.42),
            card_w - Cm(1.96),
            card_h - Cm(0.72),
            card_points,
            size=14.0,
            color=TEXT_DARK,
        )


def _build_cover_slide(slide, prs, company_name, slide_data, total, logo_path, dynamic_subtitle, cover_context, cover_bg_path):
    if cover_bg_path and os.path.exists(cover_bg_path):
        _add_full_slide_background(slide, prs, cover_bg_path)
    else:
        _draw_cover_inspired_background(slide, prs)
    _add_logo(slide, logo_path)

    title = "Strategic Growth Plan"
    prospect_name = company_name or "Prospect Name"
    month_year = slide_data.get("month_year") or os.environ.get("PPT_COVER_MONTH", "").strip()
    if not month_year:
        month_year = os.environ.get("PROSPECT_DATE", "").strip()
    if not month_year:
        month_year = "Month Year"
    month_year = _normalize_to_month_year(month_year)

    title_color = RGBColor(28, 53, 111)
    month_color = RGBColor(45, 104, 222)
    body_color = RGBColor(107, 114, 128)
    divider_color = RGBColor(176, 191, 221)

    cover_left = Inches(1.40)
    title_top = Inches(1.74)
    title_box_w = Cm(17.0)
    title_box_h = Cm(2.15)
    title_size = 50
    divider_top = Inches(3.04)
    company_top = Inches(3.20)
    company_w = Inches(6.1)
    company_h = Inches(0.5)
    month_top = Inches(3.70)
    prepared_top = Inches(4.08)

    _cover_textbox(
        slide,
        cover_left,
        title_top,
        title_box_w,
        title_box_h,
        title,
        size=title_size,
        bold=True,
        color=title_color,
        align=PP_ALIGN.LEFT,
        font_name="Calibri",
    )

    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cover_left, divider_top, Inches(0.95), Inches(0.03))
    divider.fill.solid()
    divider.fill.fore_color.rgb = divider_color
    divider.line.fill.background()

    prospect_size = 26
    if len(prospect_name) > 22:
        prospect_size = 24
    if len(prospect_name) > 30:
        prospect_size = 22
    _cover_textbox(
        slide,
        cover_left,
        company_top,
        company_w,
        company_h,
        prospect_name,
        size=prospect_size,
        bold=False,
        color=title_color,
        align=PP_ALIGN.LEFT,
        font_name="Calibri",
    )

    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.41), Inches(3.69), Inches(0.72), Inches(0.015))
    underline.fill.solid()
    underline.fill.fore_color.rgb = RGBColor(205, 63, 63)
    underline.line.fill.background()

    _cover_textbox(
        slide,
        cover_left,
        month_top,
        Inches(3.8),
        Inches(0.34),
        month_year,
        size=16,
        bold=False,
        color=month_color,
        align=PP_ALIGN.LEFT,
        font_name="Calibri",
    )
    _cover_textbox(
        slide,
        cover_left,
        prepared_top,
        Inches(4.6),
        Inches(0.34),
        "Prepared By: Traffic Radius",
        size=16,
        bold=False,
        color=body_color,
        align=PP_ALIGN.LEFT,
        font_name="Calibri",
    )


def _build_content_slide(slide, prs, session_dir, index, total, slide_data, logo_path, market_metrics=None, content_bg_path="", recent_modes=None):
    _set_bg(slide, prs)
    _add_header(slide, index, total, logo_path)

    title = _clean_optional_text(slide_data.get("title", ""))
    subtitle = _clean_optional_text(slide_data.get("subtitle", ""))
    points = _collect_existing_slide_points(slide_data)
    mode = _select_chart_mode(slide_data, index, recent_modes=recent_modes)
    if index in (4, 14):
        mode = "progress"
    chart_points = _prepare_chart_points(_parse_visual_points(slide_data.get("visual_data", []), max_points=6), mode)

    if _is_roadmap_slide(slide_data):
        _draw_roadmap_slide(slide, slide_data)
        return "roadmap"

    _add_title_block(slide, title, subtitle)

    if index == 2 and market_metrics:
        _add_slide2_kpi_panel(slide)
        _draw_slide2_company_metrics(slide, market_metrics)
        _add_bottom_bullets(slide, points, index, title, subtitle)
        return "kpi-summary"
    _add_chart_panel(slide)
    if not chart_points:
        chart_points = _prepare_chart_points(_parse_visual_points(["Visibility: 68", "Authority: 54", "Traffic: 72", "Engagement: 61", "Conversion: 58"], max_points=6), mode)
    _draw_center_chart(slide, mode, chart_points)
    if index in (5, 7, 8) and mode != "line":
        _add_line_summary_banner(slide, chart_points)

    # Insight cards below chart.
    _add_bottom_bullets(slide, points, index, title, subtitle)
    return mode


def create_ppt_package(session_dir, company_name):
    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    scripts_dir = os.path.join(project_root, "scripts")
    logo_path = _resolve_logo_path(project_root)
    cover_bg_path = _resolve_cover_background_path(project_root)
    content_bg_path = _resolve_content_background_path(project_root)
    cover_context = _load_cover_context(session_dir)
    market_metrics = _load_market_metrics(session_dir)
    dynamic_subtitle = _build_dynamic_subtitle(cover_context)

    if not session_dir or not os.path.isdir(session_dir):
        raise FileNotFoundError(f"Session directory not found: {session_dir}")

    deliv_dir = os.path.join(session_dir, "deliverables")
    os.makedirs(deliv_dir, exist_ok=True)
    output_path = os.path.join(deliv_dir, "Master_Presentation.pptx")

    data_path = os.path.join(session_dir, "ppt_slides_data.json")
    if not os.path.exists(data_path):
        print(" [1/2] Synthesizing Strategic Slide Content...")
        subprocess.run([sys.executable, os.path.join(scripts_dir, "synthesize_ppt_json.py"), session_dir, company_name], check=True)

    if not os.path.exists(data_path):
        raise FileNotFoundError(f"PPT slide data file missing: {data_path}")

    with open(data_path, "r", encoding="utf-8-sig") as f:
        slides = json.load(f)

    if not isinstance(slides, list) or not slides:
        raise ValueError("ppt_slides_data.json is empty or invalid.")

    print(" [2/2] Building editable PPTX...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    total = len(slides)
    recent_modes = []
    for i, slide_data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if i == 1:
            _build_cover_slide(slide, prs, company_name, slide_data, total, logo_path, dynamic_subtitle, cover_context, cover_bg_path)
        else:
            used_mode = _build_content_slide(
                slide,
                prs,
                session_dir,
                i,
                total,
                slide_data,
                logo_path,
                market_metrics=market_metrics,
                content_bg_path=content_bg_path,
                recent_modes=recent_modes,
            )
            if used_mode:
                recent_modes.append(used_mode)
        _resize_test_box_5(slide)

    prs.save(output_path)
    if not os.path.exists(output_path):
        raise FileNotFoundError(f"PowerPoint engine failed to save file at {output_path}")
    print(f" [+] Success: Editable PPT finalized at {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python create_presentation_pptx.py <session_dir> <company_name>")
        sys.exit(1)

    s_dir = sys.argv[1]
    c_name = sys.argv[2]

    try:
        create_ppt_package(s_dir, c_name)
    except Exception as e:
        print(f" [!] Master PPT Engine Failed: {e}")
        sys.exit(1)
